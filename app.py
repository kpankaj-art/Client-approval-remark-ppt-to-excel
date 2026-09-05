import streamlit as st
import pandas as pd
from pptx import Presentation
import io
import re

st.set_page_config(page_title="Multi-Key PPT to Excel Matcher", page_icon="🔗", layout="wide")

st.title("🔗 Dynamic Multi-Key PPT & Excel Matcher")
st.write("Analyze PPT and Excel columns dynamically, set up to 4 custom matching criteria, and export an updated Excel file.")

# File Uploaders
uploaded_excel = st.file_uploader(
    "1. Upload Excel / CSV File (.xlsx, .xls, .xlsm, .xlsb, .csv)", 
    type=["xlsx", "xls", "xlsm", "xlsb", "csv"]
)
uploaded_ppt = st.file_uploader("2. Upload PPT File (.pptx)", type=["pptx"])

def load_excel_file(file):
    filename = file.name.lower()
    
    if filename.endswith('.csv'):
        df_raw = pd.read_csv(file, nrows=15, header=None)
        file.seek(0)
    elif filename.endswith('.xlsb'):
        df_raw = pd.read_excel(file, engine='pyxlsb', nrows=15, header=None)
        file.seek(0)
    else:
        df_raw = pd.read_excel(file, nrows=15, header=None)
        file.seek(0)

    header_row_idx = 0
    max_valid_cols = 0

    for i in range(len(df_raw)):
        row_vals = [str(x).strip() for x in df_raw.iloc[i].dropna().values]
        header_candidates = [x for x in row_vals if not x.isdigit() and len(x) > 1 and x.lower() != 'nan']
        if len(header_candidates) > max_valid_cols:
            max_valid_cols = len(header_candidates)
            header_row_idx = i

    if filename.endswith('.csv'):
        df = pd.read_csv(file, header=header_row_idx)
    elif filename.endswith('.xlsb'):
        df = pd.read_excel(file, engine='pyxlsb', header=header_row_idx)
    else:
        df = pd.read_excel(file, header=header_row_idx)

    df.columns = [str(col).strip() for col in df.columns]
    df = df.loc[:, ~df.columns.str.contains('^Unnamed', na=False)]
    return df

def extract_numbers(text):
    numbers = re.findall(r'\b\d{10}\b', str(text))
    return numbers[0] if numbers else ""

def clean_key_val(val):
    if pd.isna(val) or str(val).lower() == 'nan':
        return ""
    val_str = str(val).strip()
    num = extract_numbers(val_str)
    if num:
        return num
    return re.sub(r'[^a-zA-Z0-9]', '', val_str).lower()

def extract_field(pattern, text):
    match = re.search(pattern, text, re.IGNORECASE)
    return match.group(1).strip() if match else ""

def parse_slide_content(full_text):
    text = " ".join(full_text.split())

    # Strictly Clean Outlet Name: Truncate everything starting from any key label
    clean_name = re.split(r'\s*(?:Address\s*:|Contact\s*No\s*:|Contact\s*:|District\s*:|Size\s*:|Media\s*Type\s*:|Remarks\s*:|Qty\s*:|s_no\s*:|SAP\s*Code\s*:|SAP\s*:)', text, flags=re.IGNORECASE)[0].strip()
    
    # Remove initial label prefixes if present
    pure_outlet_name = re.sub(r'^(Outlet Name|Dealer Name|Shop Name|Party Name)\s*:\s*', '', clean_name, flags=re.IGNORECASE).strip()

    # Extract Individual Fields safely
    contact_no = extract_numbers(text)
    address = extract_field(r'Address\s*:\s*(.*?)(?=\s*(?:Contact|District|Size|Media|Remarks|Qty|s_no|SAP|$))', text)
    district = extract_field(r'District\s*:\s*([A-Za-z0-9\s\-_]+?)(?=\s*(?:Size|Media|Remarks|Qty|s_no|Contact|Address|SAP|$))', text)
    size = extract_field(r'Size\s*:\s*([0-9X\s]+?)(?=\s*(?:Media|Remarks|Qty|s_no|District|Contact|Address|SAP|$))', text)
    media_type = extract_field(r'Media\s*Type\s*:\s*([A-Za-z0-9\s\-_]+?)(?=\s*(?:Remarks|Qty|s_no|District|Size|Contact|Address|SAP|$))', text)
    sap_code = extract_field(r'SAP\s*(?:Code)?\s*:\s*([A-Za-z0-9]+?)(?=\s*(?:Address|Contact|District|Size|Media|Remarks|$))', text)

    return pure_outlet_name, address, contact_no, district, size, media_type, sap_code

def process_ppt_data(ppt_file):
    prs = Presentation(ppt_file)
    ppt_records = []

    standard_keywords = [
        "address", "contact", "district", "size", "media type", "sap",
        "remarks", "qty", "s_no", "outlet name", "dealer name", "shop name"
    ]

    for idx, slide in enumerate(prs.slides):
        extra_tags = []
        text_blocks = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue

                text_lower = text.lower()
                text_blocks.append(text)

                is_standard = any(kw in text_lower for kw in standard_keywords)
                if not is_standard and len(text) <= 35:
                    extra_tags.append(text)

        full_text = " ".join(text_blocks)
        name, addr, contact, dist, sz, media, sap = parse_slide_content(full_text)

        ppt_records.append({
            "Slide_No": idx + 1,
            "PPT_Outlet_Name": name,
            "PPT_Address": addr,
            "PPT_Contact": contact,
            "PPT_District": dist,
            "PPT_Size": sz,
            "PPT_Media_Type": media,
            "PPT_SAP_Code": sap,
            "PPT_Status": " | ".join(extra_tags) if extra_tags else "Pending/None"
        })

    return pd.DataFrame(ppt_records)

def auto_select_index(options, keywords):
    for idx, opt in enumerate(options):
        opt_clean = str(opt).lower().replace("_", " ").replace(".", " ").strip()
        for kw in keywords:
            if kw in opt_clean:
                return idx
    return 0

# Execution Flow
if uploaded_excel and uploaded_ppt:
    try:
        df_excel = load_excel_file(uploaded_excel)
        df_ppt = process_ppt_data(uploaded_ppt)

        st.markdown("---")
        st.subheader("🔍 Analyzed Data Preview")
        
        tab1, tab2 = st.tabs(["📄 Uploaded Excel Preview", "🖼️ Extracted PPT Preview"])
        with tab1:
            st.dataframe(df_excel.head(10), use_container_width=True)
        with tab2:
            st.dataframe(df_ppt.head(10), use_container_width=True)

        st.markdown("---")
        st.subheader("🎯 Configure Multi-Key Matching Criteria (Up to 4 Fields)")

        excel_columns = ["-- Ignore --"] + list(df_excel.columns)
        ppt_columns = ["-- Ignore --"] + list(df_ppt.columns)

        def_ex_1 = auto_select_index(excel_columns, ["dealer / name", "dealer name", "outlet name", "name"])
        def_ex_2 = auto_select_index(excel_columns, ["dealer/ contact", "dealer contact", "contact", "mobile"])
        def_ex_3 = auto_select_index(excel_columns, ["w", "width", "size"])

        col_a, col_b, col_c, col_d = st.columns(4)

        with col_a:
            st.markdown("**Key 1 (Primary)**")
            ex_key1 = st.selectbox("Excel Col 1", options=excel_columns, index=def_ex_1, key="ex1")
            ppt_key1 = st.selectbox("PPT Col 1", options=ppt_columns, index=ppt_columns.index("PPT_Outlet_Name") if "PPT_Outlet_Name" in ppt_columns else 0, key="ppt1")

        with col_b:
            st.markdown("**Key 2**")
            ex_key2 = st.selectbox("Excel Col 2", options=excel_columns, index=def_ex_2, key="ex2")
            ppt_key2 = st.selectbox("PPT Col 2", options=ppt_columns, index=ppt_columns.index("PPT_Contact") if "PPT_Contact" in ppt_columns else 0, key="ppt2")

        with col_c:
            st.markdown("**Key 3**")
            ex_key3 = st.selectbox("Excel Col 3", options=excel_columns, index=def_ex_3, key="ex3")
            ppt_key3 = st.selectbox("PPT Col 3", options=ppt_columns, index=ppt_columns.index("PPT_Size") if "PPT_Size" in ppt_columns else 0, key="ppt3")

        with col_d:
            st.markdown("**Key 4**")
            ex_key4 = st.selectbox("Excel Col 4", options=excel_columns, key="ex4")
            ppt_key4 = st.selectbox("PPT Col 4", options=ppt_columns, key="ppt4")

        mapping_rules = []
        if ex_key1 != "-- Ignore --" and ppt_key1 != "-- Ignore --":
            mapping_rules.append((ex_key1, ppt_key1))
        if ex_key2 != "-- Ignore --" and ppt_key2 != "-- Ignore --":
            mapping_rules.append((ex_key2, ppt_key2))
        if ex_key3 != "-- Ignore --" and ppt_key3 != "-- Ignore --":
            mapping_rules.append((ex_key3, ppt_key3))
        if ex_key4 != "-- Ignore --" and ppt_key4 != "-- Ignore --":
            mapping_rules.append((ex_key4, ppt_key4))

        if not mapping_rules:
            st.warning("⚠️ Please select at least Key 1 to perform matching.")
        else:
            if st.button("▶️ Process & Append PPT Data into Excel", type="primary", use_container_width=True):
                with st.spinner("Creating composite keys and matching records..."):
                    
                    excel_key_series = pd.Series([""] * len(df_excel), index=df_excel.index)
                    for ex_col, _ in mapping_rules:
                        excel_key_series += df_excel[ex_col].apply(clean_key_val) + "_"
                    df_excel['composite_match_key'] = excel_key_series

                    ppt_key_series = pd.Series([""] * len(df_ppt), index=df_ppt.index)
                    for _, ppt_col in mapping_rules:
                        ppt_key_series += df_ppt[ppt_col].apply(clean_key_val) + "_"
                    df_ppt['composite_match_key'] = ppt_key_series

                    ppt_target_cols = [c for c in df_ppt.columns if c not in ['composite_match_key']]
                    
                    for target in ppt_target_cols:
                        mapping_dict = dict(zip(df_ppt['composite_match_key'], df_ppt[target]))
                        new_col_name = f"Matched_{target}" if target in df_excel.columns else target
                        df_excel[new_col_name] = df_excel['composite_match_key'].map(mapping_dict).fillna("Not Found / No Match" if target == "PPT_Status" else "")

                    df_final = df_excel.drop(columns=['composite_match_key'])

                    st.success("✅ Matching successfully completed!")
                    st.subheader("📋 Updated Excel Preview (With New PPT Columns)")
                    st.dataframe(df_final, use_container_width=True)

                    output = io.BytesIO()
                    with pd.ExcelWriter(output, engine='openpyxl') as writer:
                        df_final.to_excel(writer, index=False)
                    processed_data = output.getvalue()

                    st.download_button(
                        label="📥 Download Updated Excel File (.xlsx)",
                        data=processed_data,
                        file_name="Updated_Excel_With_PPT_Data.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        use_container_width=True
                    )

    except Exception as e:
        st.error(f"Error processing request: {e}")
